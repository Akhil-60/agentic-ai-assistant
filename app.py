"""
=====================================================================================
 SpaceVerse AI  —  Multi-Agent AI Assistant (Akhil's Assistant)
=====================================================================================
 A Gradio-based agentic AI app powered by Groq (LLaMA 3.3 70B) with tools, RAG,
 long-term memory, multi-agent routing, and multi-modal (text/image/audio) support.

 FILE STRUCTURE (use Ctrl+F on the section banners to jump around):
   1. IMPORTS
   2. CONFIG & CLIENT
   3. PROMPTS            (system prompt + specialist agent prompts)
   4. TOOL FUNCTIONS
        4.1 Basic utilities
        4.2 Web, search & research
        4.3 Finance & weather
        4.4 Media (image / audio / vision / TTS)
        4.5 YouTube & code execution
        4.6 Documents (create / convert / presentation)
        4.7 Resume & job tools
        4.8 Productivity (email / calendar / meeting / github)
        4.9 Planning & multi-agent helpers
   5. RAG SYSTEM         (FAISS index + document indexing)
   6. LONG-TERM MEMORY
   7. MULTI-AGENT ROUTER
   8. TOOL REGISTRY & SCHEMA
   9. MAIN AGENT LOOP    (agent_respond)
   10. GRADIO UI / LAUNCH
=====================================================================================
"""

# =====================================================================================
# 1. IMPORTS
# =====================================================================================

# ---- Standard library ----
import os
import re
import io
import json
import time
import base64
import pickle
import asyncio
import tempfile
import threading
import contextlib
import urllib.parse
from datetime import datetime
from zoneinfo import ZoneInfo

# ---- Third-party (core) ----
import requests
import numpy as np
import pandas as pd
import faiss
import gradio as gr
import sympy as sp
import fitz  # PyMuPDF
from bs4 import BeautifulSoup
from groq import Groq
from pypdf import PdfReader
from docx import Document
from fpdf import FPDF
from pptx import Presentation
from sentence_transformers import SentenceTransformer
from duckduckgo_search import DDGS
from youtube_transcript_api import YouTubeTranscriptApi
import edge_tts

# ---- Optional / heavy packages (app still runs if these are missing) ----
try:
    from github import Github
except Exception:
    Github = None
try:
    from playwright.sync_api import sync_playwright
except Exception:
    sync_playwright = None
try:
    import yagmail
except Exception:
    yagmail = None


# =====================================================================================
# 2. CONFIG & CLIENT
# =====================================================================================
# NOTE: Never hardcode API keys. Set them as environment variables before running:
#   GROQ_API_KEY, GITHUB_TOKEN, EMAIL_ID, EMAIL_PASSWORD
# If an old key was ever leaked, REVOKE it immediately.

client = Groq(api_key=os.getenv("GROQ_API_KEY"))
MODEL = "llama-3.3-70b-versatile"
DEBUG = False
MAX_TOOL_TURNS = 29


# =====================================================================================
# 3. PROMPTS
# =====================================================================================

# ---- 3.1 Main system prompt ----
SYSTEM_PROMPT = """
You are 'Akhil's Assistant' - a capable, proactive AI agent that solves tasks end-to-end, accurately and clearly.

# CORE PRINCIPLES
- Understand the user's REAL intent first, not just surface keywords.
- Be action-oriented: if a request is clear, do it fully in one go.
- Ask a clarifying question ONLY when genuinely ambiguous. At most ONE short question, then proceed.
- Reliable over impressive: a correct, grounded answer beats a confident wrong one.

# REASONING & TOOL USE
- Use a TOOL only when it genuinely helps (current info, computation, files). For things you already know (full forms, definitions, general knowledge), just answer directly - DO NOT call a tool.
- After a tool runs, weave its output into a natural answer. Never paste raw tool dumps unless asked.
- If a tool fails or returns nothing useful, say so honestly and try another approach - don't invent an answer.

# OUTPUT QUALITY
- Lead with the direct answer, then supporting detail if needed.
- Simple/factual questions (full form of X, a definition, a small fact) -> answer DIRECTLY in chat as text.
- Create a file/PPT/document ONLY when the user explicitly asks for one. If the user says "chat me batao", do NOT make any file.

# LANGUAGE & TONE
- Always reply in the SAME language the user uses (Hindi -> Hindi, Hinglish -> Hinglish, English -> English).
- Warm, clear, concise. No filler, no hype.

# CONTEXT AWARENESS (translate / explain / summarize "this")
- If the user says "translate", "translate this", "translate in <language>", "explain this", "summarize this",
  "shorten this", or "elaborate this" WITHOUT providing new text, use the most recent assistant response as the input.
- Support translation into any language the user requests.

# HONESTY
- Never invent facts, links, sources, or data. If unsure, say so plainly.
"""

# ---- 3.2 Specialist (multi-agent) prompts ----
AGENT_PROMPTS = {
    "coding": """
You are a Senior Software Engineer.
Write clean, optimized, secure, and production-ready code.
Debug errors, explain logic, suggest improvements,
and provide complete runnable solutions.
""",
    "math": """
You are an Expert Mathematician.
Solve problems step-by-step.
Show formulas and calculations clearly.
Verify answers before responding.
""",
    "research": """
You are a Professional Research Analyst.
Gather information from multiple sources.
Compare findings, identify trends,
and provide detailed reports with references.
""",
    "writing": """
You are an Expert Content Writer.
Create professional, well-structured,
clear, and engaging content.
Adapt writing style according to user requirements.
""",
    "analyst": """
You are a Data Analytics Specialist.
Analyze data, identify patterns,
generate insights, and provide recommendations.
Present findings clearly and professionally.
""",
    "meeting": """
You are a Meeting Intelligence Agent.
Generate meeting summaries,
key decisions, action items,
and follow-up recommendations.
""",
    "email": """
You are a Professional Email Assistant.
Write clear, concise,
professional, and effective emails.
Improve tone, grammar, and clarity.
""",
    "planner": """
You are a Strategic Planning Agent.
Break complex goals into actionable steps.
Create roadmaps, execution plans,
timelines, and milestones.
""",
    "reviewer": """
You are a Quality Assurance Reviewer.
Check responses for:
- Mistakes
- Missing information
- Logical issues
- Hallucinations
- Weak explanations

Improve and refine the final answer.
""",
    "executor": """
You are an Autonomous Execution Agent.
Create plans, organize workflows,
and execute tasks systematically.
Focus on practical completion of objectives.
""",
    "team": """
You are a Multi-Agent Coordinator.
Assign tasks to specialist agents,
combine outputs,
resolve conflicts,
and generate the best final response.
""",
    "job": """
You are a Fake Job Detection Specialist.

Responsibilities:
- Analyze job postings
- Detect scam patterns
- Evaluate hiring credibility
- Verify company legitimacy
- Assess salary realism
- Generate risk scores

Check for:
- Registration fees
- Upfront payments
- Telegram recruitment
- WhatsApp-only hiring
- Gmail/Yahoo email addresses
- Unrealistic salary offers
- Missing company details
- Fake urgency tactics
- Suspicious interview processes

Output Format:

Risk Score: XX/100

Risk Level:
Low / Medium / High

Positive Indicators:
- ...

Red Flags:
- ...

Final Verdict:
Safe / Suspicious / Likely Scam

Recommended Action:
...
""",
    "autopilot": """
You are SpaceVerse Autopilot - the default autonomous agent.

Understand the user's real objective first.
Automatically decide which specialist skills, tools, and steps are needed - NEVER ask the user which tool to use.
Make a short internal plan, then act: call the right tools in sequence and combine everything into ONE clear, final answer.

You can draw on these specialist styles as needed:
- coding, math, research, writing, analyst
- planner, executor, reviewer
- email, meeting, team, job (fake-job detection)

Pick the best combination automatically. Reply in the user's language.
""",
}


# =====================================================================================
# 4. TOOL FUNCTIONS
# =====================================================================================

# -------------------------------------------------------------------------------------
# 4.1 Basic utilities
# -------------------------------------------------------------------------------------
def calculator(expression):
    return str(sp.sympify(expression))


def get_current_time():
    return datetime.now(ZoneInfo("Asia/Kolkata")).strftime("%Y-%m-%d %H:%M:%S")


# -------------------------------------------------------------------------------------
# 4.2 Web, search & research
# -------------------------------------------------------------------------------------
def web_search(query):
    results = DDGS().text(query, max_results=5)
    if not results:
        return "Koi result nahi mila."
    out = []
    for r in results:
        out.append(f"- {r.get('title','')}\n  URL: {r.get('href','')}\n  {r.get('body','')}")
    return "\n\n".join(out)


def read_webpage(url):
    try:
        resp = requests.get(url, headers={"User-Agent": "Mozilla/5.0"}, timeout=10)
        soup = BeautifulSoup(resp.text, "html.parser")
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()
        return soup.get_text(separator=" ", strip=True)[:3000]
    except Exception as e:
        return f"Page nahi padh paaya: {e}"


def deep_research(topic):
    queries = [topic, f"{topic} latest", f"{topic} explained in detail"]
    seen, picked = set(), []
    for q in queries:
        try:
            results = DDGS().text(q, max_results=4)
        except Exception:
            results = []
        for r in results:
            url = r.get("href", "")
            if url and url not in seen:
                seen.add(url)
                picked.append((r.get("title", ""), url, r.get("body", "")))
        if len(picked) >= 5:
            break
    if not picked:
        return "Koi research result nahi mila."
    briefs = []
    for title, url, snippet in picked[:4]:
        page = read_webpage(url)
        page = page[:900] if isinstance(page, str) else ""
        briefs.append(f"SOURCE: {title}\nURL: {url}\n{snippet}\nPAGE: {page}")
    return "\n\n---\n\n".join(briefs)


def advanced_research(query):
    research = deep_research(query)
    prompt = f"""
    Research Results:

    {research}

    Create:
    1. Executive Summary
    2. Key Findings
    3. Risks
    4. Recommendations
    """
    response = client.chat.completions.create(
        model=MODEL, messages=[{"role": "user", "content": prompt}])
    return response.choices[0].message.content


def open_website(url):
    if sync_playwright is None:
        return "Playwright install nahi hai."
    try:
        with sync_playwright() as p:
            browser = p.chromium.launch(headless=True)
            page = browser.new_page()
            page.goto(url)
            title = page.title()
            browser.close()
            return title
    except Exception as e:
        return str(e)


def browser_search(url):
    if sync_playwright is None:
        return "Playwright install nahi hai."
    try:
        with sync_playwright() as p:
            browser = p.chromium.launch(headless=True)
            page = browser.new_page()
            page.goto(url)
            title = page.title()
            content = page.content()
            browser.close()
            return f"Title: {title}\n\n{content[:3000]}"
    except Exception as e:
        return str(e)


# -------------------------------------------------------------------------------------
# 4.3 Finance & weather
# -------------------------------------------------------------------------------------
def get_weather(city):
    try:
        geo = requests.get("https://geocoding-api.open-meteo.com/v1/search",
                           params={"name": city, "count": 1}, timeout=10).json()
        if not geo.get("results"):
            return f"'{city}' nahi mila."
        loc = geo["results"][0]
        w = requests.get("https://api.open-meteo.com/v1/forecast",
                         params={"latitude": loc["latitude"], "longitude": loc["longitude"],
                                 "current": "temperature_2m,relative_humidity_2m,wind_speed_10m"},
                         timeout=10).json()
        c = w.get("current", {})
        return (f"{loc.get('name')}, {loc.get('country','')}: "
                f"{c.get('temperature_2m')}C, humidity {c.get('relative_humidity_2m')}%, "
                f"wind {c.get('wind_speed_10m')} km/h")
    except Exception as e:
        return f"Mausam nahi mila: {e}"


def convert_currency(amount, from_currency, to_currency):
    try:
        r = requests.get("https://api.frankfurter.app/latest",
                         params={"amount": amount, "from": from_currency.upper(),
                                 "to": to_currency.upper()}, timeout=10).json()
        rates = r.get("rates", {})
        to = to_currency.upper()
        if not rates or to not in rates:
            return "Conversion nahi ho payi (currency code check karo)."
        return f"{amount} {from_currency.upper()} = {rates[to]} {to}"
    except Exception as e:
        return f"Currency error: {e}"


# -------------------------------------------------------------------------------------
# 4.4 Media (image / audio / vision / TTS)
# -------------------------------------------------------------------------------------
def generate_image(prompt):
    try:
        clean = urllib.parse.quote(prompt.strip())
        url = f"https://image.pollinations.ai/prompt/{clean}?width=1024&height=1024&nologo=true&model=flux"
        resp = requests.get(url, timeout=60)
        resp.raise_for_status()
        path = tempfile.NamedTemporaryFile(delete=False, suffix=".png").name
        with open(path, "wb") as f:
            f.write(resp.content)
        return path
    except Exception as e:
        return f"Image generate nahi ho payi: {e}"


def analyze_image(filepath, question):
    with open(filepath, "rb") as f:
        b64 = base64.b64encode(f.read()).decode()
    response = client.chat.completions.create(
        model="meta-llama/llama-4-scout-17b-16e-instruct",
        messages=[{
            "role": "user",
            "content": [
                {"type": "text", "text": question or "Is image mein kya hai? Detail mein batao."},
                {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{b64}"}},
            ],
        }],
        max_tokens=1024,
    )
    return response.choices[0].message.content


def transcribe_audio(filepath):
    with open(filepath, "rb") as f:
        result = client.audio.transcriptions.create(model="whisper-large-v3-turbo", file=f)
    return result.text


def text_to_speech(text, voice="hi-IN-SwaraNeural"):
    clean = (text or "").strip()
    if not clean:
        return None
    path = tempfile.NamedTemporaryFile(delete=False, suffix=".mp3").name
    err = {}

    def _run():
        try:
            loop = asyncio.new_event_loop()
            asyncio.set_event_loop(loop)
            loop.run_until_complete(edge_tts.Communicate(clean[:1000], voice).save(path))
            loop.close()
        except Exception as e:
            err["e"] = e

    t = threading.Thread(target=_run)
    t.start()
    t.join()
    if err:
        print("TTS error:", err["e"])
        return None
    return path


def voice_chat(audio_file):
    text = transcribe_audio(audio_file)
    response = client.chat.completions.create(
        model=MODEL, messages=[{"role": "user", "content": text}])
    answer = response.choices[0].message.content
    audio = text_to_speech(answer)
    return {"transcript": text, "answer": answer, "audio": audio}


# -------------------------------------------------------------------------------------
# 4.5 YouTube & code execution
# -------------------------------------------------------------------------------------
def youtube_summary(url):
    try:
        if "v=" in url:
            vid = url.split("v=")[1].split("&")[0]
        elif "youtu.be/" in url:
            vid = url.split("youtu.be/")[1].split("?")[0]
        else:
            vid = url.strip()
        try:
            data = YouTubeTranscriptApi().fetch(vid)
            text = " ".join(s.text for s in data)
        except Exception:
            data = YouTubeTranscriptApi.get_transcript(vid)
            text = " ".join(s["text"] for s in data)
        return text[:4000]
    except Exception as e:
        return f"Transcript nahi mila: {e}"


def run_python(code):
    output = io.StringIO()
    try:
        with contextlib.redirect_stdout(output):
            exec(code, {})
        result = output.getvalue()
        return result if result else "(Code chal gaya. Output ke liye print() use karo.)"
    except Exception as e:
        return f"Error: {e}"


# -------------------------------------------------------------------------------------
# 4.6 Documents (create / convert / presentation)
# -------------------------------------------------------------------------------------
def create_document(content, filetype="pdf", filename="document"):
    try:
        ftype = (filetype or "pdf").lower().lstrip(".")
        safe = re.sub(r"[^\w\-]", "_", str(filename)) or "document"
        content = str(content)
        if ftype in ("doc", "docx", "word"):
            path = os.path.join(tempfile.gettempdir(), f"{safe}.docx")
            d = Document()
            for line in content.split("\n"):
                d.add_paragraph(line)
            d.save(path)
        elif ftype in ("xlsx", "xls", "excel"):
            path = os.path.join(tempfile.gettempdir(), f"{safe}.xlsx")
            rows = [line.split(",") for line in content.strip().split("\n")]
            pd.DataFrame(rows).to_excel(path, index=False, header=False)
        else:
            path = os.path.join(tempfile.gettempdir(), f"{safe}.pdf")
            pdf = FPDF()
            pdf.add_page()
            pdf.set_auto_page_break(auto=True, margin=15)
            pdf.set_font("Helvetica", size=12)
            for line in content.split("\n"):
                safe_line = line.encode("latin-1", "replace").decode("latin-1")
                pdf.multi_cell(0, 8, safe_line if safe_line else " ")
            pdf.output(path)
        return path
    except Exception as e:
        return f"Document banane me error: {e}"


def convert_last_document(filetype="docx", filename="converted"):
    if not last_doc_text.strip():
        return "Koi document upload nahi hua jise convert kar saku."
    return create_document(last_doc_text, filetype, filename)


def create_presentation(slides_json, filename="presentation"):
    try:
        data = slides_json
        if isinstance(data, str):
            data = json.loads(data)
        prs = Presentation()
        for s in data:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = str(s.get("title", ""))
            body = slide.placeholders[1].text_frame
            body.clear()
            bullets = s.get("bullets", [])
            if isinstance(bullets, str):
                bullets = [bullets]
            for i, b in enumerate(bullets):
                p = body.paragraphs[0] if i == 0 else body.add_paragraph()
                p.text = str(b)
        safe = re.sub(r"[^\w\-]", "_", str(filename)) or "presentation"
        path = os.path.join(tempfile.gettempdir(), f"{safe}.pptx")
        prs.save(path)
        return path
    except Exception as e:
        return f"Presentation banane me error: {e}"


# -------------------------------------------------------------------------------------
# 4.7 Resume & job tools
# -------------------------------------------------------------------------------------
def screen_resume(job_description):
    if not last_doc_text.strip():
        return "Pehle ek resume (PDF/docx) upload karo, phir screen karunga."
    prompt = f"""You are an expert recruiter. Screen this resume against the job description.

JOB DESCRIPTION:
{job_description}

RESUME:
{last_doc_text[:6000]}

Give a structured screening:
1. Match score (0-100) with a one-line reason
2. Matching skills/experience (bullets)
3. Missing or weak areas (bullets)
4. Verdict: Strong fit / Possible fit / Not a fit
5. 3 interview questions to probe the gaps
Reply in the user's language (Hindi/Hinglish if they used it)."""
    try:
        r = client.chat.completions.create(
            model=MODEL, messages=[{"role": "user", "content": prompt}],
            temperature=0.3, max_tokens=900)
        return r.choices[0].message.content
    except Exception as e:
        return f"Screening error: {e}"


def fake_job_detector(job_text):
    prompt = f"""
You are an expert job fraud detection analyst.

Analyze this job posting:

{job_text}

Check:

1. Scam Score (0-100)
2. Risk Level (Low/Medium/High)
3. Warning Signs
4. Positive Signs
5. Final Verdict

Look for:
- Registration fees
- Telegram hiring
- WhatsApp only hiring
- Gmail/Yahoo email
- Unrealistic salary
- No company information
- Fake urgency

Return structured report.
"""
    response = client.chat.completions.create(
        model=MODEL,
        messages=[{"role": "user", "content": prompt}],
        temperature=0.2)
    return response.choices[0].message.content


# -------------------------------------------------------------------------------------
# 4.8 Productivity (email / calendar / meeting / github)
# -------------------------------------------------------------------------------------
def send_email(to, subject, body):
    try:
        yag = yagmail.SMTP(os.getenv("EMAIL_ID"), os.getenv("EMAIL_PASSWORD"))
        yag.send(to=to, subject=subject, contents=body)
        return "Email sent successfully."
    except Exception as e:
        return str(e)


def add_event(title, date):
    with open("calendar_events.txt", "a", encoding="utf-8") as f:
        f.write(f"{date} | {title}\n")
    return "Event added."


def summarize_meeting(text):
    prompt = f"""
    Meeting Transcript:

    {text}

    Generate:
    1. Summary
    2. Key Decisions
    3. Action Items
    """
    response = client.chat.completions.create(
        model=MODEL, messages=[{"role": "user", "content": prompt}])
    return response.choices[0].message.content


def github_repo_info(repo_name):
    if Github is None:
        return "PyGithub install nahi hai."
    try:
        g = Github(os.getenv("GITHUB_TOKEN"))
        repo = g.get_repo(repo_name)
        return (f"Repository: {repo.full_name}\nStars: {repo.stargazers_count}\n"
                f"Forks: {repo.forks_count}\nLanguage: {repo.language}\n"
                f"Description:\n{repo.description}")
    except Exception as e:
        return str(e)


# -------------------------------------------------------------------------------------
# 4.9 Planning & multi-agent helpers
# -------------------------------------------------------------------------------------
def create_plan(task):
    prompt = f"""
    Break this task into steps:

    {task}

    Return numbered steps only.
    """
    response = client.chat.completions.create(
        model=MODEL, messages=[{"role": "user", "content": prompt}])
    return response.choices[0].message.content


def execute_workflow(task):
    plan = create_plan(task)
    return f"Task:\n{task}\n\nPlan:\n{plan}"


def autonomous_task(task):
    plan = create_plan(task)
    return f"Task:\n{task}\n\nExecution Plan:\n{plan}"


def review_answer(question, answer):
    prompt = f"""
    User Question:

    {question}

    AI Answer:

    {answer}

    Check for mistakes, missing information, hallucinations. Improve the answer.
    """
    response = client.chat.completions.create(
        model=MODEL, messages=[{"role": "user", "content": prompt}])
    return response.choices[0].message.content


def team_agent(task):
    plan = create_plan(task)
    research = deep_research(task)
    prompt = f"""You are a coordinator. Plan aur research ka use karke task ka ek saaf, final, well-structured jawab likho. User ki language me reply karo.

TASK: {task}

PLAN:
{plan}

RESEARCH:
{research[:3000]}
"""
    try:
        r = client.chat.completions.create(
            model=MODEL, messages=[{"role": "user", "content": prompt}],
            temperature=0.4, max_tokens=900)
        return r.choices[0].message.content
    except Exception as e:
        return f"Team agent error: {e}"


# =====================================================================================
# 5. RAG SYSTEM  (FAISS: multi-doc + persistent)
# =====================================================================================

embedder = SentenceTransformer("all-MiniLM-L6-v2")
EMBED_DIM = 384
RAG_INDEX_PATH = "rag.index"
RAG_CHUNKS_PATH = "rag_chunks.pkl"

# ---- Session/runtime state (shared globals) ----
indexed_docs = []      # names of currently indexed documents
last_doc_text = ""     # full text of the most recently uploaded document


def chunk_text(text, size=500, overlap=100):
    chunks = []
    i = 0
    while i < len(text):
        chunks.append(text[i:i + size])
        i += size - overlap
    return chunks


# ---- Load existing index from disk, or start fresh ----
if os.path.exists(RAG_INDEX_PATH) and os.path.exists(RAG_CHUNKS_PATH):
    faiss_index = faiss.read_index(RAG_INDEX_PATH)
    with open(RAG_CHUNKS_PATH, "rb") as f:
        doc_chunks = pickle.load(f)
else:
    faiss_index = faiss.IndexFlatIP(EMBED_DIM)
    doc_chunks = []


def reset_rag():
    global faiss_index, doc_chunks
    faiss_index = faiss.IndexFlatIP(EMBED_DIM)
    doc_chunks = []
    indexed_docs.clear()
    for p in (RAG_INDEX_PATH, RAG_CHUNKS_PATH):
        if os.path.exists(p):
            os.remove(p)


def _add_to_rag(text):
    global doc_chunks
    chunks = [c for c in chunk_text(text) if c.strip()]
    if not chunks:
        return
    emb = np.array(embedder.encode(chunks), dtype="float32")
    faiss.normalize_L2(emb)
    faiss_index.add(emb)
    doc_chunks.extend(chunks)
    faiss.write_index(faiss_index, RAG_INDEX_PATH)
    with open(RAG_CHUNKS_PATH, "wb") as f:
        pickle.dump(doc_chunks, f)


# ---- Document readers / indexers ----
def read_scanned_pdf(filepath, max_pages=5):
    out = []
    doc = fitz.open(filepath)
    for i, page in enumerate(doc):
        if i >= max_pages:
            break
        pix = page.get_pixmap(dpi=150)
        img_path = tempfile.NamedTemporaryFile(delete=False, suffix=".png").name
        pix.save(img_path)
        try:
            out.append(analyze_image(img_path, "Is page ka saara text jaisa likha hai waise nikaalo."))
        except Exception as e:
            out.append(f"(page {i+1} padhne me error: {e})")
    doc.close()
    return "\n\n".join(out)


def index_pdf(filepath):
    global last_doc_text
    reader = PdfReader(filepath)
    text = "".join((p.extract_text() or "") + "\n" for p in reader.pages)
    if len(text.strip()) < 50:
        text = read_scanned_pdf(filepath)
    last_doc_text = text
    _add_to_rag(text)


def index_docx(filepath):
    global last_doc_text
    text = "\n".join(p.text for p in Document(filepath).paragraphs)
    last_doc_text = text
    _add_to_rag(text)


def index_excel(filepath):
    global last_doc_text
    text = pd.read_excel(filepath).to_string()
    last_doc_text = text
    _add_to_rag(text)


def analyze_excel(filepath):
    global last_doc_text
    df = pd.read_excel(filepath)
    text = df.to_string()
    last_doc_text = text
    _add_to_rag(text)
    return {"rows": len(df), "columns": len(df.columns),
            "column_names": list(df.columns), "status": "Added to RAG"}


# ---- Retrieval ----
def retrieve(query, k=5):
    if faiss_index.ntotal == 0 or not query:
        return ""
    q = np.array(embedder.encode([query]), dtype="float32")
    faiss.normalize_L2(q)
    k = min(k, faiss_index.ntotal)
    scores, idx = faiss_index.search(q, k)
    return "\n\n".join(doc_chunks[i] for i in idx[0] if 0 <= i < len(doc_chunks))


def knowledge_search(query):
    docs = retrieve(query)
    if docs:
        return docs
    return "No relevant knowledge found."


# =====================================================================================
# 6. LONG-TERM MEMORY
# =====================================================================================

MEMORY_FILE = "memory.json"

if os.path.exists(MEMORY_FILE):
    with open(MEMORY_FILE, "r", encoding="utf-8") as f:
        user_memory = json.load(f)
else:
    user_memory = {"name": "", "interests": [], "projects": [], "preferences": [], "facts": []}

# Ensure all expected keys exist (back-compat for older memory files)
for _k, _default in (("name", ""), ("interests", []), ("projects", []),
                     ("preferences", []), ("facts", [])):
    if _k not in user_memory:
        user_memory[_k] = _default


def _save_memory():
    with open(MEMORY_FILE, "w", encoding="utf-8") as f:
        json.dump(user_memory, f, ensure_ascii=False, indent=2)


def search_memory(query):
    facts = user_memory.get("facts", [])
    hits = [f for f in facts if (query or "").lower() in f.lower()]
    return "\n".join(hits) if hits else "Memory me kuch related nahi mila."


def smart_memory_search(query):
    query = (query or "").lower()
    results = []
    for key, value in user_memory.items():
        if query in str(value).lower():
            results.append(f"{key}: {value}")
    if results:
        return "\n".join(results)
    return "Memory me kuch nahi mila."


# =====================================================================================
# 7. MULTI-AGENT ROUTER
# =====================================================================================
def route_agent(user_text):
    txt = (user_text or "").lower()
    if any(x in txt for x in ["python", "code", "program", "bug", "coding", "algorithm"]):
        return "coding"
    if any(x in txt for x in ["research", "analyze", "compare", "study", "report", "investigate"]):
        return "research"
    if any(x in txt for x in ["math", "solve", "equation", "calculate", "ganit", "integral"]):
        return "math"
    if any(x in txt for x in ["teach", "explain", "tutorial", "learn"]):
        return "planner"
    if any(x in txt for x in ["email", "essay", "article", "blog"]):
        return "writing"
    if any(x in txt for x in ["excel", "csv", "dashboard", "analytics", "data"]):
        return "analyst"
    if any(x in txt for x in ["meeting", "minutes", "transcript"]):
        return "meeting"
    if any(x in txt for x in ["mail"]):
        return "email"
    if any(x in txt for x in ["calendar", "schedule", "reminder"]):
        return "planner"
    if any(x in txt for x in ["automation", "workflow", "execute", "complete task"]):
        return "executor"
    if any(x in txt for x in ["team", "multiple agents", "collaborate"]):
        return "team"
    if any(x in txt for x in ["job scam", "fake job", "job verification", "verify job", "job fraud"]):
        return "job"
    # No specific agent matched -> use Autopilot (autonomous default)
    return "autopilot"


# =====================================================================================
# 8. TOOL REGISTRY & SCHEMA
# =====================================================================================

# ---- 8.1 Name -> function map (used to actually run a tool) ----
TOOL_FUNCTIONS = {
    "calculator": calculator,
    "get_current_time": get_current_time,
    "web_search": web_search,
    "get_weather": get_weather,
    "convert_currency": convert_currency,
    "read_webpage": read_webpage,
    "youtube_summary": youtube_summary,
    "run_python": run_python,
    "create_document": create_document,
    "transcribe_audio": transcribe_audio,
    "analyze_image": analyze_image,
    "generate_image": generate_image,
    "convert_last_document": convert_last_document,
    "deep_research": deep_research,
    "create_presentation": create_presentation,
    "screen_resume": screen_resume,
    "analyze_excel": analyze_excel,
    "github_repo_info": github_repo_info,
    "browser_search": browser_search,
    "summarize_meeting": summarize_meeting,
    "send_email": send_email,
    "add_event": add_event,
    "create_plan": create_plan,
    "review_answer": review_answer,
    "search_memory": search_memory,
    "autonomous_task": autonomous_task,
    "advanced_research": advanced_research,
    "open_website": open_website,
    "execute_workflow": execute_workflow,
    "voice_chat": voice_chat,
    "smart_memory_search": smart_memory_search,
    "team_agent": team_agent,
    "knowledge_search": knowledge_search,
    "fake_job_detector": fake_job_detector,
}

# ---- 8.2 Tool schema sent to the model (descriptions guide when to call each) ----
tools = [
    {"type": "function", "function": {
        "name": "calculator",
        "description": "Exact arithmetic / math expressions ke liye (e.g. '23*17', '2**10').",
        "parameters": {"type": "object", "properties": {
            "expression": {"type": "string", "description": "Math expression jaise '23*17'"}}, "required": ["expression"]}}},

    {"type": "function", "function": {
        "name": "get_current_time",
        "description": "Abhi ka current date aur time (IST). Jab user 'aaj', 'abhi', date ya time pooche.",
        "parameters": {"type": "object", "properties": {}}}},

    {"type": "function", "function": {
        "name": "run_python",
        "description": "Python code chala ke uska ASLI output deta hai. Code me print() zaroori hai.",
        "parameters": {"type": "object", "properties": {
            "code": {"type": "string", "description": "Chalane wala Python code"}}, "required": ["code"]}}},

    {"type": "function", "function": {
        "name": "web_search",
        "description": "Internet pe TAAZA jaankari search. SIRF current/badalti info ke liye. General knowledge ke liye seedha jawab do.",
        "parameters": {"type": "object", "properties": {
            "query": {"type": "string", "description": "Saaf search keywords"}}, "required": ["query"]}}},

    {"type": "function", "function": {
        "name": "get_weather",
        "description": "Kisi sheher ka abhi ka mausam. SIRF tab jab user explicitly weather/mausam pooche.",
        "parameters": {"type": "object", "properties": {
            "city": {"type": "string", "description": "Sheher ka naam"}}, "required": ["city"]}}},

    {"type": "function", "function": {
        "name": "convert_currency",
        "description": "Ek currency se doosri me convert karta hai live rate se.",
        "parameters": {"type": "object", "properties": {
            "amount": {"type": "number", "description": "Kitni raqam"},
            "from_currency": {"type": "string", "description": "Source code jaise USD"},
            "to_currency": {"type": "string", "description": "Target code jaise INR"}},
            "required": ["amount", "from_currency", "to_currency"]}}},

    {"type": "function", "function": {
        "name": "youtube_summary",
        "description": "YouTube video ke URL se transcript laata hai.",
        "parameters": {"type": "object", "properties": {
            "url": {"type": "string", "description": "YouTube video URL"}}, "required": ["url"]}}},

    {"type": "function", "function": {
        "name": "read_webpage",
        "description": "Kisi diye gaye URL ka text content padhta hai.",
        "parameters": {"type": "object", "properties": {
            "url": {"type": "string", "description": "Webpage URL"}}, "required": ["url"]}}},

    {"type": "function", "function": {
        "name": "create_document",
        "description": "Diye content se file banata hai: PDF, Word(docx), ya Excel(xlsx). SIRF jab user explicitly file maange.",
        "parameters": {"type": "object", "properties": {
            "content": {"type": "string", "description": "File me jo likhna hai"},
            "filetype": {"type": "string", "description": "pdf, docx, ya xlsx"},
            "filename": {"type": "string", "description": "File ka naam (bina extension)"}},
            "required": ["content"]}}},

    {"type": "function", "function": {
        "name": "analyze_image",
        "description": "Image file ko dekh ke jawab deta hai.",
        "parameters": {"type": "object", "properties": {
            "filepath": {"type": "string", "description": "Image file path"},
            "question": {"type": "string", "description": "Image ke baare me sawaal"}}, "required": ["filepath"]}}},

    {"type": "function", "function": {
        "name": "transcribe_audio",
        "description": "Audio file ko text me badalta hai.",
        "parameters": {"type": "object", "properties": {
            "filepath": {"type": "string", "description": "Audio file path"}}, "required": ["filepath"]}}},

    {"type": "function", "function": {
        "name": "generate_image",
        "description": "Text se image generate karta hai. SIRF jab user explicitly image/tasveer maange.",
        "parameters": {"type": "object", "properties": {
            "prompt": {"type": "string", "description": "Image ka description"}}, "required": ["prompt"]}}},

    {"type": "function", "function": {
        "name": "deep_research",
        "description": "Kisi topic pe gehri research - kai web searches + top pages padhta hai. Jab user detail/research maange.",
        "parameters": {"type": "object", "properties": {
            "topic": {"type": "string", "description": "Research ka topic"}}, "required": ["topic"]}}},

    {"type": "function", "function": {
        "name": "convert_last_document",
        "description": "Aakhri upload ki gayi file ko doosre format me convert karta hai.",
        "parameters": {"type": "object", "properties": {
            "filetype": {"type": "string", "description": "Target: docx, pdf, ya xlsx"},
            "filename": {"type": "string", "description": "File ka naam"}},
            "required": ["filetype"]}}},

    {"type": "function", "function": {
        "name": "create_presentation",
        "description": "PowerPoint (.pptx) file banata hai. SIRF tab jab user SAAF-SAAF 'PPT'/'slides'/'presentation' maange.",
        "parameters": {"type": "object", "properties": {
            "slides_json": {"type": "string", "description": "JSON: [{\"title\":\"...\",\"bullets\":[\"p1\",\"p2\"]}, ...]"},
            "filename": {"type": "string", "description": "File ka naam"}},
            "required": ["slides_json"]}}},

    {"type": "function", "function": {
        "name": "analyze_excel",
        "description": "Excel file analyze karke insights deta hai.",
        "parameters": {"type": "object", "properties": {
            "filepath": {"type": "string"}}, "required": ["filepath"]}}},

    {"type": "function", "function": {
        "name": "github_repo_info",
        "description": "GitHub repository ki info deta hai.",
        "parameters": {"type": "object", "properties": {
            "repo_name": {"type": "string"}}, "required": ["repo_name"]}}},

    {"type": "function", "function": {
        "name": "browser_search",
        "description": "Website kholta hai aur info nikalta hai.",
        "parameters": {"type": "object", "properties": {
            "url": {"type": "string"}}, "required": ["url"]}}},

    {"type": "function", "function": {
        "name": "team_agent",
        "description": "Multiple agents ko saath use karta hai.",
        "parameters": {"type": "object", "properties": {
            "task": {"type": "string"}}, "required": ["task"]}}},

    {"type": "function", "function": {
        "name": "send_email",
        "description": "Email bhejta hai.",
        "parameters": {"type": "object", "properties": {
            "to": {"type": "string"}, "subject": {"type": "string"}, "body": {"type": "string"}},
            "required": ["to", "subject", "body"]}}},

    {"type": "function", "function": {
        "name": "create_plan",
        "description": "Step by step plan banata hai.",
        "parameters": {"type": "object", "properties": {
            "task": {"type": "string"}}, "required": ["task"]}}},

    {"type": "function", "function": {
        "name": "advanced_research",
        "description": "Advanced research karta hai.",
        "parameters": {"type": "object", "properties": {
            "query": {"type": "string"}}, "required": ["query"]}}},

    {"type": "function", "function": {
        "name": "execute_workflow",
        "description": "Multi-step workflow banata aur execute karta hai.",
        "parameters": {"type": "object", "properties": {
            "task": {"type": "string"}}, "required": ["task"]}}},

    {"type": "function", "function": {
        "name": "screen_resume",
        "description": "Aakhri upload kiye resume ko job description ke against screen karta hai.",
        "parameters": {"type": "object", "properties": {
            "job_description": {"type": "string", "description": "Job description ya role requirements"}},
            "required": ["job_description"]}}},

    {"type": "function", "function": {
        "name": "knowledge_search",
        "description": "Upload kiye gaye documents (RAG) me se relevant content dhoondta hai. Jab user ne PDF/docx/xlsx upload kiya ho aur usi ke baare me poochhe.",
        "parameters": {"type": "object", "properties": {
            "query": {"type": "string", "description": "Documents me kya dhoondna hai"}}, "required": ["query"]}}},

    {"type": "function", "function": {
        "name": "smart_memory_search",
        "description": "User ki long-term memory (name, interests, projects, preferences, facts) me search. Jab pichhli yaad rakhi gayi baat chahiye.",
        "parameters": {"type": "object", "properties": {
            "query": {"type": "string", "description": "Memory me kya dhoondna hai"}}, "required": ["query"]}}},

    {"type": "function", "function": {
        "name": "summarize_meeting",
        "description": "Meeting transcript ka summary, key decisions aur action items banata hai.",
        "parameters": {"type": "object", "properties": {
            "text": {"type": "string", "description": "Meeting ka transcript text"}}, "required": ["text"]}}},

    {"type": "function", "function": {
        "name": "fake_job_detector",
        "description": "Analyze job postings and detect scams or fake jobs.",
        "parameters": {"type": "object", "properties": {
            "job_text": {"type": "string", "description": "Job description text"}},
            "required": ["job_text"]}}},

    {"type": "function", "function": {
        "name": "add_event",
        "description": "Calendar event save karta hai. Jab user koi reminder/event/schedule add karne ko bole.",
        "parameters": {"type": "object", "properties": {
            "title": {"type": "string", "description": "Event ka naam"},
            "date": {"type": "string", "description": "Date/time, jaise '2026-06-20 17:00'"}},
            "required": ["title", "date"]}}},
]


# =====================================================================================
# 9. MAIN AGENT LOOP
# =====================================================================================
def agent_respond(message, history):
    # ---- 9.1 Parse incoming message (text + files) ----
    if isinstance(message, dict):
        user_text = message.get("text", "")
        files = message.get("files", [])
    else:
        user_text, files = message, []

    # ---- 9.2 Special commands ----
    # Clear memory
    if (user_text or "").strip().lower() == "clear memory":
        user_memory.clear()
        user_memory.update({"name": "", "interests": [], "projects": [], "preferences": [], "facts": []})
        _save_memory()
        yield "Memory cleared."
        return

    # Clear documents (RAG)
    if (user_text or "").strip().lower() in ("clear docs", "reset docs", "naye docs",
                                             "documents clear", "folder clear"):
        reset_rag()
        yield "Saare documents hata diye. Ab naya folder/PDF upload karo."
        return

    # ---- 9.3 Translation shortcut ----
    # Agar user "translate"/"translate in <lang>" bole BINA naya text diye,
    # to last assistant response ko translate karne ke liye use karo.
    translation_keywords = [
        "translate", "translation", "translate into", "translate to",
        "hindi me", "english me", "spanish me", "french me", "german me",
        "chinese me", "japanese me", "arabic me", "russian me", "korean me", "urdu me",
    ]
    if user_text and any(k in user_text.lower() for k in translation_keywords):
        last_assistant = ""
        for h in reversed(history):
            if isinstance(h, dict) and h.get("role") == "assistant":
                last_assistant = h.get("content", "")
                break
        if last_assistant:
            user_text = f"""Translate the following text according to the user's request.

Previous Assistant Response:
{last_assistant}

Translation Request:
{user_text}

Return only the translated text."""

    # ---- 9.4 Handle uploaded files ----
    image_path = None
    transcript = None
    for f in files:
        low_f = str(f).lower()
        if low_f.endswith(".pdf"):
            index_pdf(f); indexed_docs.append(os.path.basename(f))
        elif low_f.endswith(".docx"):
            index_docx(f); indexed_docs.append(os.path.basename(f))
        elif low_f.endswith(".xlsx"):
            index_excel(f); indexed_docs.append(os.path.basename(f))
        elif low_f.endswith((".wav", ".mp3", ".m4a", ".ogg", ".webm")):
            transcript = transcribe_audio(f)
            user_text = transcript
        elif low_f.endswith((".png", ".jpg", ".jpeg", ".webp", ".gif")):
            image_path = f

    # ---- 9.5 Image upload -> vision answer ----
    if image_path:
        answer = analyze_image(image_path, user_text)
        partial = ""
        for word in answer.split(" "):
            partial += word + " "
            yield partial
            time.sleep(0.02)
        audio_path = text_to_speech(answer)
        if audio_path:
            yield [
                {"role": "assistant", "content": answer},
                {"role": "assistant", "content": {"path": audio_path}},
            ]
        return

    # ---- 9.6 Text -> generate a new image ----
    low = (user_text or "").lower()
    image_words = ["image", "picture", "photo", "tasveer", "wallpaper", "logo", "drawing"]
    make_words = ["banao", "bana do", "bana de", "generate", "create", "make", "draw", "design", "chahiye"]
    if user_text and any(w in low for w in image_words) and any(w in low for w in make_words):
        yield "_(image bana raha hoon... thoda ruko)_"
        img_path = generate_image(user_text)
        if isinstance(img_path, str) and os.path.exists(img_path):
            yield [
                {"role": "assistant", "content": "Ye rahi aapki image:"},
                {"role": "assistant", "content": {"path": img_path}},
            ]
        else:
            yield f"Image generate nahi ho payi. {img_path}"
        return

    # ---- 9.7 Build the message list for the model ----
    messages = [{"role": "system", "content": SYSTEM_PROMPT}]
    if transcript:
        messages.append({"role": "system", "content": f"Audio Transcript:\n{transcript}"})
    # Only send the last 8 messages of history (keeps each request small = fewer tokens)
    for h in history[-8:]:
        if isinstance(h, dict) and isinstance(h.get("content"), str) and h["content"].strip():
            messages.append({"role": h["role"], "content": h["content"]})

    # Multi-agent specialist prompt (safe .get so unknown key se crash na ho)
    agent_type = route_agent(user_text)
    sp_prompt = AGENT_PROMPTS.get(agent_type)
    if sp_prompt:
        messages.append({"role": "system", "content": sp_prompt})

    # RAG context
    context = retrieve(user_text)
    if context:
        messages.append({"role": "system",
                         "content": f"Document se related hissa (isi se jawab dena):\n{context}"})
    if indexed_docs:
        messages.append({"role": "system",
                         "content": f"Loaded documents ({len(indexed_docs)}): " + ", ".join(indexed_docs[-20:])})

    # ---- 9.8 Auto memory extraction ----
    if "my name is" in low:
        try:
            user_memory["name"] = user_text.split("is", 1)[1].strip()
            _save_memory()
        except Exception:
            pass
    if "i love" in low:
        try:
            interest = user_text.split("love", 1)[1].strip()
            if interest and interest not in user_memory["interests"]:
                user_memory["interests"].append(interest)
                _save_memory()
        except Exception:
            pass
    if "i am working on" in low:
        try:
            project = user_text.split("on", 1)[1].strip()
            if project and project not in user_memory["projects"]:
                user_memory["projects"].append(project)
                _save_memory()
        except Exception:
            pass
    if "remember" in low:
        fact = user_text.split("remember", 1)[-1].strip(" :,-.\u0964")
        if fact and fact not in user_memory["facts"]:
            user_memory["facts"].append(fact)
            _save_memory()

    profile = (
        "User Profile:\n"
        f"Name: {user_memory.get('name','')}\n"
        f"Interests: {', '.join(user_memory.get('interests', []))}\n"
        f"Projects: {', '.join(user_memory.get('projects', []))}\n"
        f"Preferences: {', '.join(user_memory.get('preferences', []))}\n"
        f"Facts: {', '.join(user_memory.get('facts', []))}"
    )
    messages.append({"role": "system",
                     "content": "Use this user memory when relevant; don't re-ask known info.\n" + profile})

    # If the message has a URL, pre-load that page's text
    urls = re.findall(r"https?://\S+", user_text or "")
    if urls:
        messages.append({"role": "system",
                         "content": f"Webpage content:\n{read_webpage(urls[0])[:3000]}"})

    messages.append({"role": "user", "content": user_text or "(file dekho)"})

    # ---- 9.9 Tool-calling loop ----
    generated_files = []

    for _ in range(MAX_TOOL_TURNS):
        # Call the model (retry a few times on tool_use_failed)
        response = None
        for attempt in range(3):
            try:
                response = client.chat.completions.create(
                    model=MODEL, messages=messages, tools=tools,
                    tool_choice="auto", temperature=0.4, max_tokens=1024)
                break
            except Exception as e:
                if "tool_use_failed" in str(e) or "Failed to call a function" in str(e):
                    continue
                yield f"API error: {e}\nThodi der baad dobara try karo."
                return

        # Fallback: try once without tools
        if response is None:
            try:
                response = client.chat.completions.create(
                    model=MODEL, messages=messages, temperature=0.5, max_tokens=1024)
            except Exception as e:
                yield f"API error: {e}\nThodi der baad dobara try karo."
                return

        msg = response.choices[0].message

        # No tool calls -> this is the final answer
        if not msg.tool_calls:
            answer = msg.content or ""
            # NOTE: Self-review (review_answer) was disabled to save tokens -
            # it made an EXTRA LLM call per response (doubling token usage).
            # Re-enable below ONLY if you have spare quota:
            #   try:
            #       answer = review_answer(user_text, answer)
            #   except Exception:
            #       pass
            # Stream the answer word-by-word
            partial = ""
            for word in answer.split(" "):
                partial += word + " "
                yield partial
                time.sleep(0.02)
            # Attach any generated files + a TTS audio reply
            out = [{"role": "assistant", "content": answer}]
            for fp in generated_files:
                out.append({"role": "assistant", "content": {"path": fp}})
            audio_path = text_to_speech(answer)
            if audio_path:
                out.append({"role": "assistant", "content": {"path": audio_path}})
            if len(out) > 1:
                yield out
            return

        # Tool calls -> run each tool, append results, then loop again
        yield "_(wait...)_"
        messages.append(msg)
        for tc in msg.tool_calls:
            fn = TOOL_FUNCTIONS.get(tc.function.name)
            if fn is None:
                result = f"Unknown tool: {tc.function.name}"
            else:
                try:
                    args = json.loads(tc.function.arguments or "{}")
                    if not isinstance(args, dict):
                        args = {}
                    result = fn(**args)
                except Exception as e:
                    result = f"Error: {e}"
            # File-producing tools: hand the file to the user
            if (tc.function.name in ("create_document", "generate_image",
                                     "convert_last_document", "create_presentation")
                    and isinstance(result, str) and os.path.exists(result)):
                generated_files.append(result)
                result = f"File ban gayi aur user ko de di: {os.path.basename(result)}"
            messages.append({"role": "tool", "tool_call_id": tc.id, "content": str(result)})

    # ---- 9.10 Loop exhausted (too many tool turns) ----
    yield "Bahut zyada tool calls ho gaye - yahin ruk raha hoon. Sawaal thoda simple karke poochho."


# =====================================================================================
# 10. GRADIO UI / LAUNCH
# =====================================================================================
if __name__ == "__main__":
    gr.ChatInterface(
        agent_respond,
        textbox=gr.MultimodalTextbox(sources=["upload", "microphone"]),
        title="🚀 SpaceVerse AI",
        description="Multi-Agent AI Platform",
    ).launch()
